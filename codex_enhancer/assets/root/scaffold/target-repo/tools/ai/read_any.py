#!/usr/bin/env python3
"""Extract useful text or structure from common local files."""

from __future__ import annotations

import argparse
import csv
import html.parser
import io
import json
import pprint
import sys
import tomllib
import zipfile
from pathlib import Path
from xml.etree import ElementTree


TEXT_SUFFIXES = {".txt", ".md", ".rst", ".py", ".js", ".ts", ".tsx", ".jsx", ".css", ".html", ".xml"}
STRUCTURED_TEXT_SUFFIXES = {".json", ".toml", ".yaml", ".yml", ".csv"}
RICH_SUFFIXES = {".pdf", ".docx", ".pptx", ".xlsx", ".odt", ".ods"}
DEPENDENCY_GROUPS = {
    "beautifulsoup4": "requirements-codex-readers.txt",
    "openpyxl": "requirements-codex-readers.txt",
    "pypdf": "requirements-codex-readers.txt",
    "pypdf or PyMuPDF": "requirements-codex-readers.txt",
    "PyMuPDF": "requirements-codex-readers.txt",
    "python-docx": "requirements-codex-readers.txt",
    "python-pptx": "requirements-codex-readers.txt",
    "ruamel.yaml": "requirements-codex-readers.txt",
}


class TextHTMLParser(html.parser.HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        stripped = data.strip()
        if stripped:
            self.parts.append(stripped)


def print_limited(text: str, *, max_chars: int) -> None:
    if len(text) <= max_chars:
        print(text)
        return
    print(text[:max_chars])
    print(f"\n[truncated at {max_chars} characters]")


def dependency_message(package: str, purpose: str) -> None:
    group_file = DEPENDENCY_GROUPS.get(package, "requirements-codex.txt")
    print(
        f"Optional dependency missing: install `{package}` from {group_file} "
        f"or the all-in requirements-codex.txt to read {purpose}."
    )


def decode_bytes(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp1252"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            pass
    try:
        from charset_normalizer import from_bytes  # type: ignore[import-not-found]
    except ImportError:
        return data.decode("utf-8", errors="replace")
    best = from_bytes(data).best()
    if best is None:
        return data.decode("utf-8", errors="replace")
    return str(best)


def read_text(path: Path) -> str:
    return decode_bytes(path.read_bytes())


def read_json(path: Path) -> str:
    data = json.loads(read_text(path))
    return json.dumps(data, indent=2, ensure_ascii=False)


def read_toml(path: Path) -> str:
    data = tomllib.loads(read_text(path))
    return pprint.pformat(data, width=100)


def read_yaml(path: Path) -> str:
    try:
        from ruamel.yaml import YAML  # type: ignore[import-not-found]
    except ImportError:
        dependency_message("ruamel.yaml", "parsed YAML; raw text follows")
        return read_text(path)
    yaml = YAML(typ="safe")
    data = yaml.load(read_text(path))
    return pprint.pformat(data, width=100)


def read_html(path: Path) -> str:
    text = read_text(path)
    try:
        from bs4 import BeautifulSoup  # type: ignore[import-not-found]
    except ImportError:
        parser = TextHTMLParser()
        parser.feed(text)
        return "\n".join(parser.parts)
    soup = BeautifulSoup(text, "html.parser")
    title = soup.title.string.strip() if soup.title and soup.title.string else None
    body = soup.get_text("\n", strip=True)
    return f"Title: {title}\n\n{body}" if title else body


def read_xml(path: Path) -> str:
    root = ElementTree.fromstring(read_text(path))
    lines = [f"Root tag: {root.tag}"]
    for element in root.iter():
        text = " ".join((element.text or "").split())
        if text:
            lines.append(f"{element.tag}: {text}")
        if len(lines) >= 200:
            lines.append("[xml element limit reached]")
            break
    return "\n".join(lines)


def read_csv(path: Path, *, max_rows: int) -> str:
    text = read_text(path)
    reader = csv.reader(io.StringIO(text))
    lines: list[str] = []
    for index, row in enumerate(reader):
        if index >= max_rows:
            lines.append(f"[truncated at {max_rows} rows]")
            break
        lines.append(" | ".join(row))
    return "\n".join(lines)


def read_pdf(path: Path) -> tuple[str, int]:
    try:
        from pypdf import PdfReader  # type: ignore[import-not-found]
    except ImportError:
        try:
            import fitz  # type: ignore[import-not-found]
        except ImportError:
            dependency_message("pypdf or PyMuPDF", "PDF files")
            return "", 1
        document = fitz.open(path)
        return "\n\n".join(page.get_text() for page in document), 0

    reader = PdfReader(str(path))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages), 0


def read_docx(path: Path) -> tuple[str, int]:
    try:
        from docx import Document  # type: ignore[import-not-found]
    except ImportError:
        dependency_message("python-docx", "DOCX files")
        return "", 1
    document = Document(str(path))
    lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            lines.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(lines), 0


def read_pptx(path: Path) -> tuple[str, int]:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError:
        dependency_message("python-pptx", "PPTX files")
        return "", 1
    deck = Presentation(str(path))
    lines: list[str] = []
    for index, slide in enumerate(deck.slides, start=1):
        lines.append(f"# Slide {index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                lines.append(text)
    return "\n".join(lines), 0


def read_xlsx(path: Path, *, max_rows: int) -> tuple[str, int]:
    try:
        from openpyxl import load_workbook  # type: ignore[import-not-found]
    except ImportError:
        dependency_message("openpyxl", "XLSX files")
        return "", 1
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        lines.append(f"# Sheet: {sheet.title}")
        for index, row in enumerate(sheet.iter_rows(values_only=True)):
            if index >= max_rows:
                lines.append(f"[truncated at {max_rows} rows]")
                break
            values = ["" if value is None else str(value) for value in row]
            lines.append(" | ".join(values).rstrip())
    return "\n".join(lines), 0


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def read_open_document(path: Path, *, max_rows: int) -> tuple[str, int]:
    try:
        with zipfile.ZipFile(path) as archive:
            content = archive.read("content.xml")
    except (KeyError, zipfile.BadZipFile, OSError) as error:
        return f"Could not read OpenDocument content.xml: {error}", 1

    root = ElementTree.fromstring(content)
    if path.suffix.lower() == ".odt":
        text_parts = [
            "".join(element.itertext()).strip()
            for element in root.iter()
            if local_name(element.tag) in {"p", "h"}
        ]
        return "\n".join(part for part in text_parts if part), 0

    lines: list[str] = []
    row_count = 0
    for row in root.iter():
        if local_name(row.tag) != "table-row":
            continue
        values: list[str] = []
        for cell in row:
            if local_name(cell.tag) == "table-cell":
                values.append(" ".join("".join(cell.itertext()).split()))
        if values:
            lines.append(" | ".join(values))
            row_count += 1
        if row_count >= max_rows:
            lines.append(f"[truncated at {max_rows} rows]")
            break
    return "\n".join(lines), 0


def extract(path: Path, *, max_chars: int, max_rows: int) -> int:
    suffix = path.suffix.lower()
    try:
        if suffix == ".json":
            print_limited(read_json(path), max_chars=max_chars)
            return 0
        if suffix == ".toml":
            print_limited(read_toml(path), max_chars=max_chars)
            return 0
        if suffix in {".yaml", ".yml"}:
            print_limited(read_yaml(path), max_chars=max_chars)
            return 0
        if suffix == ".html":
            print_limited(read_html(path), max_chars=max_chars)
            return 0
        if suffix == ".xml":
            print_limited(read_xml(path), max_chars=max_chars)
            return 0
        if suffix == ".csv":
            print_limited(read_csv(path, max_rows=max_rows), max_chars=max_chars)
            return 0
        if suffix == ".pdf":
            text, status = read_pdf(path)
            print_limited(text, max_chars=max_chars)
            return status
        if suffix == ".docx":
            text, status = read_docx(path)
            print_limited(text, max_chars=max_chars)
            return status
        if suffix == ".pptx":
            text, status = read_pptx(path)
            print_limited(text, max_chars=max_chars)
            return status
        if suffix == ".xlsx":
            text, status = read_xlsx(path, max_rows=max_rows)
            print_limited(text, max_chars=max_chars)
            return status
        if suffix in {".odt", ".ods"}:
            text, status = read_open_document(path, max_rows=max_rows)
            print_limited(text, max_chars=max_chars)
            return status
        if suffix in TEXT_SUFFIXES or suffix in STRUCTURED_TEXT_SUFFIXES or not suffix:
            print_limited(read_text(path), max_chars=max_chars)
            return 0
    except Exception as error:
        print(f"Could not read {path}: {error}")
        return 1

    print(f"Unsupported file type: {suffix or '<none>'}. Supported rich formats: {', '.join(sorted(RICH_SUFFIXES))}")
    return 2


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("path", help="file to read")
    parser.add_argument("--max-chars", type=int, default=20000, help="maximum characters to print")
    parser.add_argument("--max-rows", type=int, default=40, help="maximum spreadsheet or CSV rows to print")
    args = parser.parse_args(argv)

    path = Path(args.path).resolve()
    if not path.exists():
        print(f"File not found: {path}")
        return 2
    if not path.is_file():
        print(f"Not a file: {path}")
        return 2
    return extract(path, max_chars=max(1, args.max_chars), max_rows=max(1, args.max_rows))


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
